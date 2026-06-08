from fastapi import UploadFile
from pathlib import Path
import asyncio
import pdfplumber
from docx import Document
from pptx import Presentation


def pdf_extractor(file: Path)->list[str]:
    text = ""
    with pdfplumber.open(file) as pdf:
        for idx, page in enumerate(pdf.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text += f"Page {idx}: {page_text}\n"
            except Exception as page_error:
                raise Exception(str(page_error))
    return text
        
def word_extractor(file: Path)->str:
    doc = Document(file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def pptx_extractor(file: Path)->str:
    prs = Presentation(file)
    text_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_content.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text_content.append(paragraph.text)
    return "\n".join(text_content)    

async def file_content_extract(file_path: Path = None)->tuple[bool, str]:
    '''
    return bool result of content extraction result,
    return content is str format
    '''
    content = ""
    try:
        if file_path.exists():
            if file_path.suffix == ".pdf":
                content = await asyncio.to_thread(pdf_extractor, file_path)
            elif file_path.suffix in [".pptx", ".ppt"]:
                content = await asyncio.to_thread(pptx_extractor, file_path)
            elif file_path.suffix in [".doc", ".docx"]:
                content = await asyncio.to_thread(word_extractor, file_path)
            else:
                raise Exception(f"File type of {file_path.name} is not supported..")
            return (True, content)
        else:
            return (False, "File does not exist")
    except Exception as e:
        return (False, str(e))
        